"""
팔레오 매장 매니저 교육용 PPT 생성기

10페이지 구성:
1. 표지
2. 로그인 / 시작하기
3. 출퇴근 체크
4. 판매 입력
5. 반품 처리
6. 회원(고객) 가입
7. 재고 요청
8. 재고 입고확인 및 조회
9. 휴무 신청
10. 마무리 / 문의처
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOT_DIR = os.path.join(OUT_DIR, 'screenshots')

# ===== 색상 (팔레오/한국스마트물류 톤) =====
ORANGE     = RGBColor(0xE6, 0x51, 0x00)  # 메인 액센트
ORANGE_LT  = RGBColor(0xFF, 0xF3, 0xE0)  # 밝은 배경
NAVY       = RGBColor(0x1A, 0x35, 0x4F)  # 진한 본문 색
NAVY_GRAD  = RGBColor(0x0D, 0x22, 0x38)  # 그라데이션 어두운 색
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LT   = RGBColor(0xE8, 0xF5, 0xE9)
BLUE       = RGBColor(0x15, 0x65, 0xC0)
BLUE_LT    = RGBColor(0xE3, 0xF2, 0xFD)
RED        = RGBColor(0xC6, 0x28, 0x28)
RED_LT     = RGBColor(0xFF, 0xEB, 0xEE)
PURPLE     = RGBColor(0x6A, 0x1B, 0x9A)
PURPLE_LT  = RGBColor(0xF3, 0xE5, 0xF5)
TEAL       = RGBColor(0x1A, 0x4D, 0x4F)
TEXT_DARK  = RGBColor(0x10, 0x10, 0x10)
TEXT_MID   = RGBColor(0x40, 0x40, 0x40)
TEXT_LIGHT = RGBColor(0x80, 0x80, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_GRAY    = RGBColor(0xF8, 0xF8, 0xF8)
BORDER     = RGBColor(0xE0, 0xE0, 0xE0)

# ===== 폰트 =====
FONT_BODY  = '맑은 고딕'
FONT_HEAD  = '맑은 고딕'

# ===== 슬라이드 크기 (와이드 16:9) =====
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_blank():
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *,
             font=FONT_BODY, size=14, bold=False,
             color=TEXT_DARK, align='left', valign='top'):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    if valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = line_width
    return shape


def add_rounded_rect(slide, x, y, w, h, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.08
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = line_width
    return shape


def add_placeholder(slide, x, y, w, h, label, hint=''):
    """스크린샷 자리. 회색 점선 박스 + 라벨"""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = BG_GRAY
    box.line.color.rgb = BORDER
    box.line.width = Pt(1.5)
    # dash style — python-pptx 지원 제한적이라 실선으로 유지
    box.shadow.inherit = False
    add_text(slide, x, y + Emu(int((h.emu - Inches(0.6).emu) / 2)), w, Inches(0.6),
             f'📸 {label}', font=FONT_BODY, size=18, bold=True,
             color=TEXT_LIGHT, align='center', valign='middle')
    if hint:
        add_text(slide, x, y + h - Inches(0.6), w, Inches(0.4),
                 hint, size=11, color=TEXT_LIGHT, align='center')


def add_image_or_placeholder(slide, x, y, w, h, image_path, label, hint=''):
    """이미지 있으면 삽입, 없으면 placeholder"""
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, x, y, w, h)
            return
        except Exception:
            pass
    add_placeholder(slide, x, y, w, h, label, hint)


def page_header(slide, page_num, title, subtitle=''):
    """슬라이드 상단 헤더 — 페이지 번호 + 제목"""
    # 좌측 페이지 번호 배지
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.35), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    badge.shadow.inherit = False
    add_text(slide, Inches(0.5), Inches(0.35), Inches(0.55), Inches(0.55),
             str(page_num), size=18, bold=True, color=WHITE, align='center', valign='middle')

    # 제목
    add_text(slide, Inches(1.2), Inches(0.32), Inches(11), Inches(0.55),
             title, size=24, bold=True, color=NAVY, valign='middle')

    # 부제
    if subtitle:
        add_text(slide, Inches(1.2), Inches(0.85), Inches(11), Inches(0.3),
                 subtitle, size=12, color=TEXT_LIGHT)

    # 하단 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.3), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = ORANGE
    line.line.fill.background(); line.shadow.inherit = False


def page_footer(slide):
    """슬라이드 하단 푸터"""
    add_text(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
             '팔레오 매장 운영 시스템 · 매니저 사용 가이드', size=9, color=TEXT_LIGHT, align='center')


def step_card(slide, x, y, w, h, step_num, step_title, step_desc, color=ORANGE):
    """단계 카드 — 번호 동그라미 + 제목 + 설명"""
    # 외곽 박스
    add_rounded_rect(slide, x, y, w, h, fill=WHITE, line=BORDER, line_width=Pt(1))
    # 번호 동그라미
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = color
    circle.line.fill.background(); circle.shadow.inherit = False
    add_text(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.5), Inches(0.5),
             str(step_num), size=16, bold=True, color=WHITE, align='center', valign='middle')
    # 제목
    add_text(slide, x + Inches(0.8), y + Inches(0.18), w - Inches(1.0), Inches(0.35),
             step_title, size=13, bold=True, color=NAVY)
    # 설명
    add_text(slide, x + Inches(0.8), y + Inches(0.55), w - Inches(1.0), h - Inches(0.7),
             step_desc, size=11, color=TEXT_MID)


def tip_box(slide, x, y, w, h, text, kind='info'):
    """팁/경고 박스"""
    palette = {
        'info':    (BLUE_LT, BLUE,    '💡'),
        'warn':    (ORANGE_LT, ORANGE, '⚠️'),
        'success': (GREEN_LT, GREEN,   '✅'),
    }
    bg, border, icon = palette.get(kind, palette['info'])
    add_rounded_rect(slide, x, y, w, h, fill=bg, line=border, line_width=Pt(1))
    add_text(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), h - Inches(0.2),
             f'{icon} {text}', size=11, color=TEXT_DARK)


def notes_box(slide, x, y, w, h, title, bullets):
    """중요사항 박스 - 제목 + 불릿 리스트 (단일 textframe으로 자동 줄바꿈)"""
    add_rounded_rect(slide, x, y, w, h, fill=ORANGE_LT, line=ORANGE, line_width=Pt(1.2))
    # 제목
    add_text(slide, x + Inches(0.18), y + Inches(0.08), w - Inches(0.36), Inches(0.32),
             f'📌 {title}', size=13, bold=True, color=ORANGE)
    # 불릿 본문 — 단일 textframe에 줄바꿈으로 모든 bullet 포함 (자동 word-wrap)
    body_text = '\n'.join(f'• {b}' for b in bullets)
    body_box = slide.shapes.add_textbox(
        x + Inches(0.25), y + Inches(0.42),
        w - Inches(0.45), h - Inches(0.5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = body_text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = ln
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = TEXT_DARK


# ════════════════════════════════════════════════════════════════════
# Slide 1: 표지
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, NAVY_GRAD)

# 우측 하단 그라데이션 액센트
accent = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4.5), Inches(8), Inches(8))
accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
accent.line.fill.background(); accent.shadow.inherit = False

# 좌측 상단 작은 마크
add_text(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.4),
         '한국스마트물류', size=11, bold=True, color=RGBColor(0xFD, 0xD8, 0x35), font=FONT_BODY)
add_text(s, Inches(0.8), Inches(1.05), Inches(8), Inches(0.4),
         'STORE MANAGEMENT SYSTEM', size=10, color=WHITE, font=FONT_BODY)

# 메인 타이틀
add_text(s, Inches(0.8), Inches(2.7), Inches(11), Inches(1.1),
         '매장 매니저 사용 가이드', size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.85), Inches(11), Inches(0.6),
         '팔레오 매장 운영 시스템', size=22, color=RGBColor(0xFD, 0xD8, 0x35))

# 7개 항목 미리보기
topics = ['1. 출퇴근 체크', '2. 판매 입력', '3. 반품 처리', '4. 회원 가입',
          '5. 재고 요청', '6. 재고 입고확인 / 조회', '7. 휴무 신청']
add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
         '교육 항목 7가지', size=12, bold=True, color=WHITE)
for i, t in enumerate(topics):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + Inches(2.9 * col)
    y = Inches(5.4) + Inches(0.4 * row)
    add_text(s, x, y, Inches(2.8), Inches(0.35), '• ' + t,
             size=11, color=WHITE)

# 하단
add_text(s, Inches(0.8), Inches(7.0), Inches(12), Inches(0.4),
         '(주)한국생활건강 — 본 교육자료는 매장 매니저 분들의 시스템 활용을 돕기 위해 제작되었습니다',
         size=10, color=RGBColor(0xCC, 0xCC, 0xCC))

# ════════════════════════════════════════════════════════════════════
# Slide 2: 시작하기 / 로그인
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 0, '시작하기 · 로그인', '시스템 접속과 첫 화면 안내')

# 좌측: 로그인 화면 캡처
login_img = os.path.join(SCREENSHOT_DIR, '01-login.png')
add_image_or_placeholder(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(4.5),
                          login_img, '로그인 화면',
                          'https://goodmd-dashboard.vercel.app/')

# 우측: 설명
add_text(s, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.5),
         '📌 로그인 방법', size=16, bold=True, color=ORANGE)

step_card(s, Inches(7.3), Inches(2.2), Inches(5.5), Inches(1.1), 1,
          '브라우저로 시스템 접속',
          'Chrome 또는 Edge 브라우저에서\n시스템 주소를 입력합니다.')

step_card(s, Inches(7.3), Inches(3.4), Inches(5.5), Inches(1.1), 2,
          '아이디 / 비밀번호 입력',
          '매장별로 발급받은 아이디와\n비밀번호를 입력합니다.')

step_card(s, Inches(7.3), Inches(4.6), Inches(5.5), Inches(1.1), 3,
          '"로그인" 클릭',
          '본인의 매장 정보가 표시되는\n홈 화면으로 진입합니다.')

tip_box(s, Inches(7.3), Inches(5.9), Inches(5.5), Inches(0.7),
        '아이디·비밀번호를 잊으셨다면 본사 담당자에게 문의하세요.', kind='info')

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 3: 1. 출퇴근 체크
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 1, '출퇴근 체크', '하루 시작과 마무리 — 출근 / 퇴근 기록')

# 좌측: 사이드바 캡처
add_image_or_placeholder(s, Inches(0.5), Inches(1.6), Inches(3.0), Inches(4.2),
                          os.path.join(SCREENSHOT_DIR, '02-sidebar.png'),
                          '사이드바 출퇴근 패널',
                          '매장 매니저 로그인 후 좌측 사이드바 상단 (홈 버튼 아래)')

# 우측: 단계
add_text(s, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.5),
         '📌 출퇴근 기록 방법', size=16, bold=True, color=ORANGE)

step_card(s, Inches(7.3), Inches(2.2), Inches(5.5), Inches(1.1), 1,
          '좌측 사이드바 확인',
          '홈 버튼 아래에 출퇴근 패널이\n표시됩니다.')

step_card(s, Inches(7.3), Inches(3.4), Inches(5.5), Inches(1.1), 2,
          '"🟢 출근" / "🔴 퇴근" 클릭',
          '아침엔 출근, 퇴근 시엔 퇴근\n버튼을 클릭합니다.')

step_card(s, Inches(7.3), Inches(4.6), Inches(5.5), Inches(1.1), 3,
          '근무자 선택 후 확인',
          '본인 이름을 선택하고\n확인 버튼을 누릅니다.')

notes_box(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0),
          '중요사항',
          ['출근 / 퇴근 모두 정확한 시간에 기록 (지각·조퇴는 자동 기록되어 급여 산정에 반영됩니다)',
           '체크 누락 시 본사 근태관리 메뉴에서 수동 수정 요청 가능 — 자주 누락되지 않도록 주의해주세요'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 4: 2. 판매 입력
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 2, '판매 입력', '고객 판매 발생 시 즉시 기록')

add_image_or_placeholder(s, Inches(0.5), Inches(1.6), Inches(7.5), Inches(4.6),
                          os.path.join(SCREENSHOT_DIR, '03-sales-input.png'),
                          '판매 입력 화면',
                          '좌측 메뉴 → 판매 입력')

# 우측 단계 (작게 6단계)
add_text(s, Inches(8.3), Inches(1.6), Inches(4.7), Inches(0.4),
         '📌 입력 순서', size=14, bold=True, color=ORANGE)

steps = [
    ('상품 선택', '브랜드 선택 후 상품명/코드로 검색'),
    ('수량 / 단가', '정상가, 할인금액 입력'),
    ('결제 방법', '카드/현금/증정/시식/적립금'),
    ('🚚 택배', '고객이 택배 요청 시 체크'),
    ('회원 적립 (선택)', '회원이면 검색하여 연결'),
    ('저장', '하단 "저장" 클릭'),
]
for i, (t, d) in enumerate(steps):
    y = Inches(2.05 + i * 0.78)
    # 번호
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), y, Inches(0.35), Inches(0.35))
    circle.fill.solid(); circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background(); circle.shadow.inherit = False
    add_text(s, Inches(8.3), y, Inches(0.35), Inches(0.35), str(i+1),
             size=11, bold=True, color=WHITE, align='center', valign='middle')
    add_text(s, Inches(8.75), y, Inches(4.25), Inches(0.35), t,
             size=12, bold=True, color=NAVY, valign='middle')
    add_text(s, Inches(8.75), y + Inches(0.35), Inches(4.25), Inches(0.3), d,
             size=10, color=TEXT_MID)

notes_box(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.65),
          '중요사항',
          ['여러 상품 한 번에 판매 시 우측 [상품추가] 버튼으로 라인 추가 · 저장 시 매장재고 자동 차감 · 적립금 사용·적립은 회원 검색 후 가능'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 5: 3. 반품 처리
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 3, '반품 처리', '판매 취소 / 부분 반품')

add_image_or_placeholder(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(4.2),
                          os.path.join(SCREENSHOT_DIR, '04-sales-return.png'),
                          '반품 접수 화면',
                          '좌측 메뉴 → 반품 접수')

add_text(s, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.5),
         '📌 반품 방법', size=16, bold=True, color=ORANGE)

step_card(s, Inches(7.3), Inches(2.2), Inches(5.5), Inches(1.1), 1,
          '"반품 접수" 메뉴 진입',
          '좌측 메뉴 → 반품 접수 클릭')

step_card(s, Inches(7.3), Inches(3.4), Inches(5.5), Inches(1.1), 2,
          '판매 내역 검색',
          '날짜 / 고객 / 상품명으로 원래\n판매 건을 찾습니다.')

step_card(s, Inches(7.3), Inches(4.6), Inches(5.5), Inches(1.1), 3,
          '반품 수량 입력 → 확정',
          '전체 / 부분 반품 모두 가능.\n수량 입력 후 확정합니다.')

notes_box(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0),
          '중요사항',
          ['반품 시 매장재고는 자동 가산되지 않음 → 필요시 본사에 별도 안내 또는 매장재고 메뉴에서 수정',
           '적립금이 사용된 판매를 반품할 경우 적립금 복구는 별도 처리 — 본사에 문의'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 6: 4. 회원(고객) 가입
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 4, '회원(고객) 가입', '두 가지 가입 방식 안내')

# 좌측 — QR 가입
add_rounded_rect(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(4.4),
                 fill=BLUE_LT, line=BLUE, line_width=Pt(1.5))
add_text(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.4),
         '📱 QR 가입 (추천)', size=16, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(0.85),
         '메뉴 → QR 가입 → 고객이 QR 스캔 → 휴대폰에서 정보 입력 → 자동 등록\n✓ 빠르고 정확  ✓ 고객이 직접 확인 (입력 실수 없음)',
         size=10, color=TEXT_DARK)
add_image_or_placeholder(s, Inches(0.7), Inches(2.95), Inches(5.8), Inches(2.85),
                          os.path.join(SCREENSHOT_DIR, '05-customer-qr.png'),
                          'QR 가입 화면', '좌측 메뉴 → QR 가입')

# 우측 — 서류 가입
add_rounded_rect(s, Inches(6.9), Inches(1.5), Inches(6.2), Inches(4.4),
                 fill=GREEN_LT, line=GREEN, line_width=Pt(1.5))
add_text(s, Inches(7.1), Inches(1.6), Inches(6), Inches(0.4),
         '📝 서류 가입', size=16, bold=True, color=GREEN)
add_text(s, Inches(7.1), Inches(2.0), Inches(5.9), Inches(0.85),
         '메뉴 → 서류 가입 → 직원이 정보 입력(이름·연락처·생년월일) → 마케팅 동의 체크 → 저장\n✓ QR 사용 어려운 고객용',
         size=10, color=TEXT_DARK)
add_image_or_placeholder(s, Inches(7.1), Inches(2.95), Inches(5.8), Inches(2.85),
                          os.path.join(SCREENSHOT_DIR, '06-customer-doc.png'),
                          '서류 가입 화면', '좌측 메뉴 → 서류 가입')

notes_box(s, Inches(0.5), Inches(5.95), Inches(12.6), Inches(1.05),
          '중요사항',
          ['마케팅 수신동의는 가입일 기준 1년 유효 — 만료 14일 전 자동 안내 SMS 발송, 무응답 시 1년 자동 연장',
           '수신거부 원하는 회원은 무료수신거부 전화(0808092009)로 안내 · 생년월일 입력 시 생일 안내 메시지 발송 가능'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 7: 5. 재고 요청
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 5, '재고 요청', '재고가 부족할 때 본사에 요청')

add_image_or_placeholder(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(4.2),
                          os.path.join(SCREENSHOT_DIR, '07-stock-request.png'),
                          '재고 요청 화면',
                          '좌측 메뉴 → 재고 관리 → 재고 요청')

add_text(s, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.5),
         '📌 재고 요청 방법', size=16, bold=True, color=ORANGE)

step_card(s, Inches(7.3), Inches(2.2), Inches(5.5), Inches(1.0), 1,
          '재고 요청 메뉴 진입',
          '재고 관리 → 재고 요청')

step_card(s, Inches(7.3), Inches(3.3), Inches(5.5), Inches(1.0), 2,
          '상품 검색 후 수량 입력',
          '브랜드/상품명 검색, 필요 수량 입력')

step_card(s, Inches(7.3), Inches(4.4), Inches(5.5), Inches(1.0), 3,
          '메모 (선택) + 저장',
          '요청 사유나 특이사항이 있으면 메모')

step_card(s, Inches(7.3), Inches(5.5), Inches(5.5), Inches(1.0), 4,
          '처리 상태 확인',
          '하단 "요청 이력"에서 대기 / 완료 확인')

notes_box(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0),
          '중요사항',
          ['본사는 센터(중앙창고)에서 보내거나 다른 매장에서 재고이동으로 보낼 수 있습니다 — 출처는 신경 쓰지 않으셔도 됩니다',
           '처리 상태가 "✅ 완료"로 바뀌면 본사 출고 진행한 것 — 며칠 내 도착, 도착 후 "발주 확인 → 입고 확인" 탭에서 처리'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 8: 6. 재고 입고확인 / 조회
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 6, '재고 입고확인 / 조회', '도착한 물품 확인 + 매장재고 보기')

# 상단: 입고확인 (4분할 좌우)
add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
         '🔵 입고 확인 — 도착한 물품을 시스템에 반영', size=15, bold=True, color=BLUE)
add_image_or_placeholder(s, Inches(0.5), Inches(2.05), Inches(7), Inches(3.5),
                          os.path.join(SCREENSHOT_DIR, '08-purchase-receive.png'),
                          '발주 확인 페이지 → "입고 확인" 탭',
                          '재고 관리 → 발주 확인 → 입고 확인 탭')

add_text(s, Inches(7.8), Inches(2.1), Inches(5.2), Inches(0.4),
         '📌 입고 처리 절차', size=13, bold=True, color=ORANGE)
add_text(s, Inches(7.8), Inches(2.55), Inches(5.2), Inches(3),
         '1. 발주 확인 메뉴 → "입고 확인" 탭\n2. 도착한 발주/이동 행 펼치기\n3. 받은 수량 입력\n4. ✅ 정상 / ❌ 이상 표시\n5. "입고확인 완료" 클릭\n\n· 본사 발주, 매장이동 모두\n  한 곳에서 처리\n· 입고 즉시 매장재고에 반영',
         size=11, color=TEXT_DARK)

# 하단: 매장재고 조회 (한 줄)
add_text(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.4),
         '🟢 매장재고 조회 — 재고 관리 → 매장재고 메뉴에서 본인 매장 모든 상품 재고 확인 / [수정]으로 실사 후 수량 조정 가능',
         size=11, color=TEXT_DARK)

notes_box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
          '중요사항',
          ['입고 확인 시 실제 받은 수량을 정확히 입력 (수량 차이 시 본사 연락) — 입고확인 즉시 매장재고 자동 가산',
           '본사 발주든 매장이동이든 같은 "입고 확인" 탭에서 처리 — 어디서 왔는지 구분 안 해도 됨'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 9: 7. 휴무 신청
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
page_header(s, 7, '휴무 신청', '매월 25일까지 다음달 휴무계획 제출')

add_image_or_placeholder(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(4.2),
                          os.path.join(SCREENSHOT_DIR, '09-leave-plan.png'),
                          '휴무 신청 화면',
                          '좌측 메뉴 → 근태 관리 → 휴무 신청')

add_text(s, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.5),
         '📌 신청 방법', size=16, bold=True, color=ORANGE)

step_card(s, Inches(7.3), Inches(2.2), Inches(5.5), Inches(1.0), 1,
          '근태 관리 → 휴무 신청',
          '좌측 메뉴에서 진입')

step_card(s, Inches(7.3), Inches(3.3), Inches(5.5), Inches(1.0), 2,
          '다음달 휴무일 선택',
          '캘린더에서 원하는 휴무 날짜 클릭')

step_card(s, Inches(7.3), Inches(4.4), Inches(5.5), Inches(1.0), 3,
          '저장',
          '하단 "제출" 버튼 클릭')

notes_box(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0),
          '중요사항',
          ['두 명 이상 근무 매장의 경우 — 먼저 한 명이 신청한 날은 다른 매니저가 선택할 수 없음 (중복 방지)',
           '매월 20일~25일 사이 제출 마감 — 미제출 시 알림 표시, 다음달 시작 전까지 변경 가능'])

page_footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 10: 마무리 / 문의
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, NAVY_GRAD)

# 우측 액센트
accent = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(4), Inches(8), Inches(8))
accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
accent.line.fill.background(); accent.shadow.inherit = False

add_text(s, Inches(0.8), Inches(0.7), Inches(8), Inches(0.4),
         '교육 완료', size=11, bold=True, color=RGBColor(0xFD, 0xD8, 0x35))
add_text(s, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
         'CONGRATULATIONS', size=10, color=WHITE)

add_text(s, Inches(0.8), Inches(2.2), Inches(11), Inches(1.5),
         '수고하셨습니다 🎉', size=42, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
         '이제 7가지 기능을 자유롭게 사용하실 수 있습니다.', size=18, color=RGBColor(0xFD, 0xD8, 0x35))

# 문의처 카드
add_rounded_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.9),
                 fill=WHITE, line=None)
add_text(s, Inches(1.1), Inches(4.7), Inches(11), Inches(0.45),
         '🆘 도움이 필요할 때', size=16, bold=True, color=ORANGE)

# 3열로 문의처
cols = [
    ('💬 본사 담당자', '시스템 사용 / 회원 가입\n재고 / 발주 문의'),
    ('🔧 기술 지원', '오류 발생 / 화면 이상\n로그인 문제'),
    ('📚 매뉴얼 다시 보기', '본 자료를 다시 확인하거나\n본사에 추가 자료 요청'),
]
for i, (title, desc) in enumerate(cols):
    x = Inches(1.1 + i * 3.8)
    add_text(s, x, Inches(5.3), Inches(3.5), Inches(0.4),
             title, size=12, bold=True, color=NAVY)
    add_text(s, x, Inches(5.7), Inches(3.5), Inches(0.6),
             desc, size=10, color=TEXT_MID)

# 하단
add_text(s, Inches(0.8), Inches(6.8), Inches(12), Inches(0.4),
         '(주)한국생활건강 · 팔레오 · 본 자료는 매장 매니저 내부 교육용입니다',
         size=9, color=RGBColor(0xCC, 0xCC, 0xCC), align='center')

# ════════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════════
out_path = os.path.join(OUT_DIR, '팔레오_매장_매니저_사용가이드.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  Slides: {len(prs.slides)}')
print(f'  Size: {os.path.getsize(out_path):,} bytes')
